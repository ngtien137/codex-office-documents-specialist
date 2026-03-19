#!/usr/bin/env python3
"""Create a PPTX deck from structured JSON."""

from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict, Iterable, List

import requests
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


FORMAT_DIMENSIONS = {
    "16:9": (13.333, 7.5),
    "4:3": (10, 7.5),
}


def read_spec(path: str | None, inline_json: str | None) -> Dict[str, Any]:
    if inline_json is not None:
        data = json.loads(inline_json)
    elif path is None:
        data = json.loads(sys.stdin.read())
    else:
        data = json.loads(Path(path).read_text(encoding="utf-8"))
    if isinstance(data, list):
        return {"slides": data}
    return data


def choose_layout(prs: Presentation, *preferred_indexes: int):
    for index in preferred_indexes:
        if 0 <= index < len(prs.slide_layouts):
            return prs.slide_layouts[index]
    return prs.slide_layouts[0]


def fit_slide_size(prs: Presentation, aspect_ratio: str) -> None:
    width, height = FORMAT_DIMENSIONS.get(aspect_ratio, FORMAT_DIMENSIONS["16:9"])
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)


def set_title(slide, title: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        return
    textbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    frame = textbox.text_frame
    frame.text = title
    frame.paragraphs[0].font.size = Pt(26)
    frame.paragraphs[0].font.bold = True


def set_subtitle(slide, subtitle: str) -> None:
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == 1:
            placeholder.text = subtitle
            return
    textbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.9))
    frame = textbox.text_frame
    frame.text = subtitle
    frame.paragraphs[0].font.size = Pt(18)


def normalized_bullets(items: Iterable[Any]) -> List[Dict[str, Any]]:
    normalized: List[Dict[str, Any]] = []
    for item in items:
        if isinstance(item, dict):
            normalized.append({"text": str(item.get("text", "")), "level": int(item.get("level", 0))})
        else:
            normalized.append({"text": str(item), "level": 0})
    return normalized


def fill_bullets(text_frame, bullets: Iterable[Any]) -> None:
    text_frame.clear()
    for index, item in enumerate(normalized_bullets(bullets)):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = item["text"]
        paragraph.level = max(0, min(4, item["level"]))
        paragraph.font.size = Pt(20 if item["level"] == 0 else 18)


def add_table_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 5, 1, 6))
    set_title(slide, spec.get("title", spec.get("slide_title", "Table")))

    rows = spec.get("rows") or spec.get("table_data") or []
    if not rows:
        return
    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    shape = slide.shapes.add_table(row_count, col_count, Inches(0.7), Inches(1.5), Inches(12), Inches(5.2))
    table = shape.table
    for row_index, row in enumerate(rows):
        padded = list(row) + [""] * (col_count - len(row))
        for col_index, value in enumerate(padded):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(16 if row_index == 0 else 14)
                if row_index == 0:
                    paragraph.font.bold = True


def add_content_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 1, 5, 6))
    set_title(slide, spec.get("title", spec.get("slide_title", "Content")))
    body = None
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx != 0 and hasattr(placeholder, "text_frame"):
            body = placeholder
            break
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.2), Inches(4.8))
    fill_bullets(body.text_frame, spec.get("bullets") or spec.get("slide_text") or [])


def add_two_column_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 5, 6))
    set_title(slide, spec.get("title", spec.get("slide_title", "Two Column")))
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(4.8))
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.8))
    fill_bullets(left_box.text_frame, spec.get("left") or spec.get("left_column") or [])
    fill_bullets(right_box.text_frame, spec.get("right") or spec.get("right_column") or [])


def download_image(url: str) -> Path:
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    suffix = Path(url).suffix or ".png"
    handle = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    handle.write(response.content)
    handle.close()
    return Path(handle.name)


def resolve_image_path(spec: Dict[str, Any], temp_paths: List[Path]) -> Path:
    image_value = spec.get("image_path") or spec.get("image") or spec.get("image_url")
    if image_value is None:
        raise ValueError("image slide requires image_path, image, or image_url")
    image_text = str(image_value)
    if image_text.startswith("http://") or image_text.startswith("https://"):
        temp_path = download_image(image_text)
        temp_paths.append(temp_path)
        return temp_path
    return Path(image_text)


def add_image_slide(prs: Presentation, spec: Dict[str, Any], temp_paths: List[Path]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 5, 6))
    title = spec.get("title") or spec.get("slide_title")
    if title:
        set_title(slide, str(title))
    image_path = resolve_image_path(spec, temp_paths)
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(11.7))
    caption = spec.get("caption") or spec.get("image_caption")
    if caption:
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5))
        frame = textbox.text_frame
        frame.text = str(caption)
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        frame.paragraphs[0].font.size = Pt(14)


def add_quote_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 6, 5))
    title = spec.get("title") or spec.get("slide_title")
    if title:
        set_title(slide, str(title))
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11.2), Inches(3.6))
    frame = textbox.text_frame
    frame.text = str(spec.get("quote") or spec.get("quote_text") or "")
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    frame.paragraphs[0].font.size = Pt(26)
    frame.paragraphs[0].font.italic = True
    author = spec.get("author") or spec.get("quote_author")
    if author:
        paragraph = frame.add_paragraph()
        paragraph.text = f"- {author}"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True


def add_title_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 0, 1, 6))
    set_title(slide, spec.get("title", spec.get("slide_title", "Title")))
    subtitle = spec.get("subtitle") or spec.get("author") or ""
    if subtitle:
        set_subtitle(slide, str(subtitle))


def add_section_slide(prs: Presentation, spec: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(choose_layout(prs, 2, 5, 6))
    set_title(slide, spec.get("title", spec.get("slide_title", "Section")))
    subtitle = spec.get("subtitle", "")
    if subtitle:
        set_subtitle(slide, str(subtitle))


def add_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = text
    except Exception:
        return


def build_presentation(spec: Dict[str, Any]) -> tuple[Presentation, Dict[str, Any], List[Path]]:
    template = spec.get("template")
    prs = Presentation(template) if template else Presentation()
    if not template:
        fit_slide_size(prs, str(spec.get("format", "16:9")))

    slide_type_counts: Dict[str, int] = {}
    temp_paths: List[Path] = []

    for slide_spec in spec.get("slides", []):
        slide_type = str(slide_spec.get("type") or slide_spec.get("slide_type") or "content").lower()
        before_count = len(prs.slides)
        if slide_type == "title":
            add_title_slide(prs, slide_spec)
        elif slide_type == "section":
            add_section_slide(prs, slide_spec)
        elif slide_type == "table":
            add_table_slide(prs, slide_spec)
        elif slide_type == "image":
            add_image_slide(prs, slide_spec, temp_paths)
        elif slide_type == "quote":
            add_quote_slide(prs, slide_spec)
        elif slide_type == "two_column":
            add_two_column_slide(prs, slide_spec)
        else:
            add_content_slide(prs, slide_spec)

        new_slide = prs.slides[len(prs.slides) - 1]
        add_speaker_notes(new_slide, str(slide_spec.get("speaker_notes") or ""))
        slide_type_counts[slide_type] = slide_type_counts.get(slide_type, 0) + 1
        if len(prs.slides) == before_count:
            raise RuntimeError(f"Failed to append slide for type: {slide_type}")

    stats = {
        "slide_count": len(prs.slides),
        "slide_type_counts": slide_type_counts,
        "used_template": bool(template),
        "format": str(spec.get("format", "16:9")),
    }
    return prs, stats, temp_paths


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create a PPTX deck from structured JSON.")
    parser.add_argument("input", nargs="?", help="Optional JSON file path. If omitted, reads stdin unless --json is used.")
    parser.add_argument("--json", help="Inline JSON specification.")
    parser.add_argument("--output", required=True, help="Output PPTX path.")
    parser.add_argument("--summary", action="store_true", help="Print a JSON summary after writing the presentation.")
    return parser.parse_args()


def cleanup_temp_paths(paths: Iterable[Path]) -> None:
    for path in paths:
        try:
            path.unlink(missing_ok=True)
        except Exception:
            continue


def main() -> int:
    args = parse_args()
    spec = read_spec(args.input, args.json)
    presentation, stats, temp_paths = build_presentation(spec)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        presentation.save(str(output_path))
    finally:
        cleanup_temp_paths(temp_paths)

    if args.summary:
        print(
            json.dumps(
                {
                    "path": str(output_path.resolve()),
                    "stats": stats,
                },
                ensure_ascii=False,
                indent=2,
            )
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
