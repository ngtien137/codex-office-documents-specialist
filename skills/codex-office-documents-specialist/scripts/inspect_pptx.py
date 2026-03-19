#!/usr/bin/env python3
"""Inspect PPTX/PPTM files and emit a JSON summary."""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def preview(text: str, limit: int = 140) -> str:
    compact = " ".join(text.split())
    if len(compact) <= limit:
        return compact
    return compact[: limit - 3].rstrip() + "..."


def iter_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def notes_text(slide) -> str:
    try:
        text = slide.notes_slide.notes_text_frame.text
        return preview(text)
    except Exception:
        return ""


def slide_summary(slide, index: int) -> Dict[str, Any]:
    title_text = slide.shapes.title.text if slide.shapes.title is not None else ""
    picture_count = 0
    table_count = 0
    chart_count = 0
    media_count = 0
    text_samples: List[str] = []

    all_shapes = list(iter_shapes(slide.shapes))
    for shape in all_shapes:
        if getattr(shape, "has_table", False):
            table_count += 1
        if getattr(shape, "has_chart", False):
            chart_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        if shape.shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT}:
            media_count += 1
        if getattr(shape, "has_text_frame", False):
            text = preview(shape.text_frame.text)
            if text:
                text_samples.append(text)

    notes = notes_text(slide)
    return {
        "slide_index": index,
        "title": preview(title_text),
        "shape_count": len(all_shapes),
        "picture_count": picture_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "media_count": media_count,
        "has_notes": bool(notes),
        "notes_preview": notes,
        "text_samples": text_samples[:6],
    }


def media_files(path: Path) -> List[str]:
    with zipfile.ZipFile(path) as archive:
        return sorted(name for name in archive.namelist() if name.startswith("ppt/media/"))


def presentation_properties(prs: Presentation) -> Dict[str, Any]:
    props = prs.core_properties
    return {
        "title": props.title,
        "subject": props.subject,
        "author": props.author,
        "keywords": props.keywords,
        "comments": props.comments,
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
        "last_modified_by": props.last_modified_by,
    }


def build_report(path: Path) -> Dict[str, Any]:
    prs = Presentation(str(path))
    slides = [slide_summary(slide, index + 1) for index, slide in enumerate(prs.slides)]
    media = media_files(path)

    return {
        "path": str(path.resolve()),
        "properties": presentation_properties(prs),
        "summary": {
            "slide_count": len(prs.slides),
            "slide_master_count": len(prs.slide_masters),
            "slide_layout_count": len(prs.slide_layouts),
            "media_file_count": len(media),
            "picture_count": sum(item["picture_count"] for item in slides),
            "table_count": sum(item["table_count"] for item in slides),
            "chart_count": sum(item["chart_count"] for item in slides),
            "slides_with_notes": sum(1 for item in slides if item["has_notes"]),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
        },
        "slides": slides,
        "media_files": media,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Inspect a PPTX or PPTM file and emit a JSON summary.")
    parser.add_argument("path", help="Path to the PPTX or PPTM file.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    report = build_report(Path(args.path))
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
